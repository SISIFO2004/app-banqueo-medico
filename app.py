import streamlit as st
import PyPDF2
from pptx import Presentation
from PIL import Image
import io
from procesador import configurar_api, extraer_calamares_y_preguntas

# Configuración de página ancha para ver Calamares y Flashcards en paralelo
st.set_page_config(page_title="Banqueo Médico Multimodal", layout="wide")

if "data_procesada" not in st.session_state:
    st.session_state.data_procesada = None

# --- FUNCIÓN DE PRE-FILTRADO DE IMÁGENES ---
def es_tabla_o_esquema(imagen):
    """
    Analiza la complejidad de la imagen.
    Las tablas tienen fondos planos y pocos colores únicos.
    Las fotos clínicas (sangre, piel, órganos) tienen miles de colores.
    """
    try:
        img_analisis = imagen.copy()
        img_analisis.thumbnail((150, 150))
        # Si tiene más de 3000 colores, es probable que sea una foto real y no una tabla
        colores = img_analisis.getcolors(3000)
        
        if colores is None:
            return False # Detectado como foto (complejo)
        return True # Detectado como tabla/esquema (plano)
    except Exception:
        return True # En caso de error, permitir para no perder info

st.title("🦑 Sistema de Banqueo y Calamares Mentales")

# --- BARRA LATERAL / CONTROLES ---
st.markdown("### ⚙️ Configuración del Banqueo")
col_ctrl1, col_ctrl2 = st.columns(2)
with col_ctrl1:
    num_preguntas = st.slider("Número de Flashcards a generar", min_value=1, max_value=50, value=30)
with col_ctrl2:
    tema_sugerido = st.text_input("Tema sugerido (Opcional)", placeholder="Ej: Tratamiento. Si está vacío, serán todos.")

tema_final = tema_sugerido if tema_sugerido.strip() != "" else "todos los temas"
st.markdown("---")

# --- CARGA DE DOCUMENTOS ---
archivo_subido = st.file_uploader("Sube tu apunte (PDF o PPTX)", type=["pdf", "pptx"])

texto_extraido = ""
imagenes_brutas = []

if archivo_subido is not None:
    nombre_archivo = archivo_subido.name.lower()
    
    # Procesar PDF
    if nombre_archivo.endswith(".pdf"):
        lector = PyPDF2.PdfReader(archivo_subido)
        for pagina in lector.pages:
            if pagina.extract_text():
                texto_extraido += pagina.extract_text() + "\n"
                
    # Procesar PPTX (Extrayendo texto e imágenes de tablas)
    elif nombre_archivo.endswith(".pptx"):
        prs = Presentation(archivo_subido)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texto_extraido += shape.text + "\n"
                
                if hasattr(shape, "image"):
                    image_bytes = shape.image.blob
                    imagen = Image.open(io.BytesIO(image_bytes))
                    if imagen.mode != 'RGB':
                        imagen = imagen.convert('RGB')
                    imagenes_brutas.append(imagen)

    # --- INTERFAZ DE FILTRADO DE IMÁGENES ---
    imagenes_a_enviar = []
    if imagenes_brutas:
        st.info("🤖 **Pre-filtrado inteligente:** Se han marcado automáticamente tablas y esquemas. Desmarca las fotos clínicas si el sistema se equivocó.")
        
        columnas = st.columns(4) 
        for idx, img in enumerate(imagenes_brutas):
            col_idx = idx % 4
            
            # Clasificación automática
            es_util = es_tabla_o_esquema(img)
            
            with columnas[col_idx]:
                st.image(img, use_container_width=True)
                incluir = st.checkbox(f"Incluir", value=es_util, key=f"img_{idx}")
                
                if incluir:
                    imagenes_a_enviar.append(img)
                if not es_util:
                    st.caption("🔴 *Detectado como Foto Clínica*")

    # --- BOTÓN DE PROCESAMIENTO ---
    if st.button("🚀 Generar Banqueo y Calamares", type="primary"):
        if not texto_extraido and not imagenes_a_enviar:
            st.error("Sube un documento válido o selecciona imágenes.")
        else:
            with st.spinner(f"Procesando {len(imagenes_a_enviar)} imágenes y creando {num_preguntas} flashcards..."):
                try:
                    # Configurar la API con el secreto de Streamlit
                    api_key = st.secrets["GEMINI_API_KEY"]
                    configurar_api(api_key)
                    
                    # Llamada al motor de procesamiento
                    resultado = extraer_calamares_y_preguntas(
                        texto_extraido, imagenes_a_enviar, num_preguntas, tema_final
                    )
                    st.session_state.data_procesada = resultado
                except Exception as e:
                    st.error(f"Error en el servidor de IA: {e}")

# --- SECCIÓN DE RESULTADOS ---
if st.session_state.data_procesada:
    data = st.session_state.data_procesada
    
    st.success(f"📌 **Tema detectado:** {data.get('tema_general', 'No especificado')}")
    st.write("") 
    
    col_izq, col_der = st.columns([1, 1])
    
    # Visualización de Calamares Mentales
    with col_izq:
        st.header("🦑 Calamares Mentales")
        for categoria, contenido in data.get("calamares", {}).items():
            if contenido:
                with st.expander(f"🔹 {categoria}", expanded=True):
                    # Manejo de seguridad por si la IA envía diccionarios en lugar de texto
                    if isinstance(contenido, dict):
                        texto_mostrar = contenido.get('contenido', '') + "\n\n" + contenido.get('mnemotecnia', '')
                    else:
                        texto_mostrar = contenido
                    
                    st.markdown(texto_mostrar)
                
    # Visualización de Flashcards
    with col_der:
        st.header("🗂️ Flashcards de Banqueo")
        for idx, card in enumerate(data.get("flashcards", [])):
            with st.expander(f"Pregunta {idx + 1}: {card.get('pregunta', '')}"):
                st.success(f"**Respuesta:** {card.get('respuesta', '')}")
